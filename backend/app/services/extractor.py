from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation


IMAGE_EXTENSIONS = {".emf", ".gif", ".jpeg", ".jpg", ".png", ".svg", ".tif", ".tiff", ".webp", ".wmf"}


def count_pptx_images(path: Path) -> int:
    try:
        with ZipFile(path) as archive:
            return sum(
                1
                for name in archive.namelist()
                if name.startswith("ppt/media/") and Path(name).suffix.lower() in IMAGE_EXTENSIONS
            )
    except Exception:
        return 0


def extract_pptx(path: Path) -> tuple[str, int]:
    try:
        deck = Presentation(str(path))
        text_parts: list[str] = []
        for index, slide in enumerate(deck.slides, start=1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"Slide {index}: " + "\n".join(slide_text))
        return "\n\n".join(text_parts), len(deck.slides)
    except Exception:
        return extract_pptx_xml(path)


def extract_pptx_xml(path: Path) -> tuple[str, int]:
    text_parts: list[str] = []
    slide_count = 0
    with ZipFile(path) as archive:
        slide_names = sorted(name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
        slide_count = len(slide_names)
        for index, name in enumerate(slide_names, start=1):
            raw = archive.read(name).decode("utf-8", errors="ignore")
            fragments: list[str] = []
            start = 0
            while True:
                start = raw.find("<a:t>", start)
                if start == -1:
                    break
                start += len("<a:t>")
                end = raw.find("</a:t>", start)
                if end == -1:
                    break
                fragments.append(raw[start:end])
                start = end
            if fragments:
                text_parts.append(f"Slide {index}: " + " ".join(fragments))
    return "\n\n".join(text_parts), slide_count


def extract_document_text(path: Path, file_type: str) -> tuple[str, int | None]:
    if file_type == "pptx":
        return extract_pptx(path)
    if file_type == "ppt":
        return "", None
    raise ValueError(f"Unsupported file type: {file_type}")
