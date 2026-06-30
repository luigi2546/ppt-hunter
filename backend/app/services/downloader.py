import hashlib
from pathlib import Path
from urllib.parse import urlparse

import httpx
from pptx import Presentation

from app.core.config import settings


MAX_DOWNLOAD_BYTES = 150 * 1024 * 1024


def download_to_storage(document_id: str, url: str, file_type: str) -> tuple[Path, str, int]:
    target = settings.storage_dir / "raw" / f"{document_id}.{file_type}"
    if urlparse(url).netloc == "example.com" and file_type == "pptx":
        return create_mock_deck(target, url)

    digest = hashlib.sha256()
    size = 0

    with httpx.stream("GET", url, follow_redirects=True, timeout=60) as response:
        response.raise_for_status()
        content_type = response.headers.get("content-type", "").lower()
        if "html" in content_type:
            raise ValueError(f"Expected PowerPoint file, got content-type {content_type}")

        with target.open("wb") as output:
            for chunk in response.iter_bytes(chunk_size=1024 * 1024):
                size += len(chunk)
                if size > MAX_DOWNLOAD_BYTES:
                    raise ValueError("File exceeds maximum download size")
                digest.update(chunk)
                output.write(chunk)

    return target, digest.hexdigest(), size


def create_mock_deck(target: Path, url: str) -> tuple[Path, str, int]:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[0])
    slide.shapes.title.text = "PPT Hunter Sample"
    slide.placeholders[1].text = f"Generated locally for {url}"

    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "AI Discovery Pipeline"
    slide.placeholders[1].text = "Search, dedupe, download, extract, categorize, summarize, and index."

    deck.save(target)
    payload = target.read_bytes()
    return target, hashlib.sha256(payload).hexdigest(), len(payload)
